"""
Generate LJMU Thesis Defence Presentation
=========================================
Professional PowerPoint for Master's thesis defence.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
COLOR_PRIMARY = RGBColor(0, 51, 102)      # Navy blue
COLOR_SECONDARY = RGBColor(0, 102, 153)   # Teal
COLOR_ACCENT = RGBColor(255, 102, 0)      # Orange
COLOR_TEXT = RGBColor(51, 51, 51)         # Dark gray
COLOR_LIGHT = RGBColor(240, 248, 255)     # Light blue bg

def add_title_slide(prs, title, subtitle, author="Raghunandan Kavi", id="PN1196933"):
    blank_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank_layout)
    
    # Background shape
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_PRIMARY
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.3), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(200, 220, 240)
    p.alignment = PP_ALIGN.CENTER
    
    # Author info
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.3), Inches(1))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{author}  |  Student ID: {id}"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Institution
    inst_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.6))
    tf = inst_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Liverpool John Moores University  |  MSc Artificial Intelligence and Machine Learning"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(180, 200, 220)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, section_num, title):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_SECONDARY
    bg.line.fill.background()
    
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"0{section_num}" if section_num < 10 else str(section_num)
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.3), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, extra_text=None):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Light background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(250, 250, 250)
    bg.line.fill.background()
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(14)
        p.level = 0
    
    if extra_text:
        p = tf.add_paragraph()
        p.text = extra_text
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_SECONDARY
        p.space_before = Pt(20)
        p.font.italic = True
    
    return slide

def add_image_slide(prs, title, image_path, caption=None):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(250, 250, 250)
    bg.line.fill.background()
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1.5), Inches(1.4), width=Inches(10.3))
    
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.6))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_SECONDARY
        p.alignment = PP_ALIGN.CENTER
        p.font.italic = True
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(250, 250, 250)
    bg.line.fill.background()
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.9), Inches(5.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECONDARY
    p.space_after = Pt(10)
    
    for item in left_items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(10)
    
    # Divider
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.55), Inches(1.4), Inches(0.02), Inches(5.3))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_SECONDARY
    line.line.fill.background()
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(5.9), Inches(5.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECONDARY
    p.space_after = Pt(10)
    
    for item in right_items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(10)
    
    return slide

# ==================== BUILD PRESENTATION ====================

# Slide 1: Title
add_title_slide(prs,
    "A Comparative Study of Retrieval-Augmented Generation Pipelines for Medical Question Answering",
    "Reducing Hallucination in Clinical Decision Support Systems")

# Slide 2: Agenda
add_content_slide(prs, "Presentation Outline", [
    "Background & Motivation — Why medical AI needs grounded answers",
    "Problem Statement — Hallucination risks in clinical settings",
    "Research Questions — What this study investigates",
    "Methodology — Five pipeline configurations under controlled conditions",
    "Results — RAGAS and DeepEval metrics across all pipelines",
    "Key Findings & Implications — What this means for clinical deployment",
    "Conclusion & Future Work — Where we go from here"
])

# Slide 3: Section — Background
add_section_slide(prs, 1, "Background & Motivation")

# Slide 4: The Problem
add_content_slide(prs, "The Hallucination Problem in Medical AI", [
    "Large language models produce confident but factually incorrect medical advice",
    "Documented case: chatbot incorrectly stated a chemotherapy drug was safe during pregnancy",
    "Hallucination rates range from 15-30% on complex medical queries (Singhal et al., 2023)",
    "In oncology, meta-analysis found 23% hallucination rate across 6,523 responses (Yoon et al., 2026)",
    "Current systems lack evidence grounding — clinicians cannot verify claims against sources"
], "Patient safety demands answers that are traceable, verifiable, and grounded in authoritative sources.")

# Slide 5: What is RAG?
add_content_slide(prs, "Retrieval-Augmented Generation (RAG)", [
    "RAG grounds LLM outputs in external knowledge sources before generating answers",
    "Three components: Retriever → Encoder → Generator (Lewis et al., 2020)",
    "Dense retrieval uses vector embeddings to find semantically similar documents",
    "Sparse retrieval uses keyword matching for exact terminology",
    "Medical documents require BOTH: semantic coverage + terminological precision"
], "The central question: which retrieval strategy most effectively reduces hallucination in medical QA?")

# Slide 6: Section — Research Questions
add_section_slide(prs, 2, "Research Questions & Objectives")

# Slide 7: RQs
add_content_slide(prs, "Research Questions", [
    "RQ1: How do the four retrieval-based pipelines differ from the vanilla baseline in reducing hallucinations?",
    "RQ2: How does retrieval depth affect faithfulness, context quality, and hallucination risk?",
    "RQ3: How do query expansion, hybrid retrieval, and query reformulation compare across different medical question types?",
    "RQ4: Where do failures mainly occur, and how do failure patterns vary by strategy?"
])

# Slide 8: Hypotheses
add_content_slide(prs, "Hypotheses", [
    "H1: All four RAG pipelines achieve significantly higher faithfulness and lower hallucination than the vanilla baseline",
    "H2: Hybrid retrieval achieves the highest faithfulness and lowest hallucination due to combining semantic and lexical search",
    "H3: Multi-query expansion achieves highest context recall but lowest precision, trading coverage for noise"
])

# Slide 9: Section — Methodology
add_section_slide(prs, 3, "Methodology")

# Slide 10: Dataset
add_content_slide(prs, "Dataset: MedQuAD", [
    "47,457 question-answer pairs from 12 NIH sources (Ben Abacha & Demner-Fushman, 2019)",
    "Covers symptoms, treatments, diagnosis, medication, prognosis, prevention",
    "Consumer health information — no personally identifiable data",
    "Preprocessing: duplicate removal, HTML stripping, whitespace normalization, ASCII conversion",
    "Stratified sampling ensures proportional representation across question types"
])

# Slide 11: Controlled Variables
add_two_column_slide(prs, "Controlled Experimental Design",
    "Fixed Across All Pipelines",
    ["LLM: GPT-4o-mini (temperature=0.0)", "Embedding: text-embedding-3-small", "Chunking: 400 tokens, 50 overlap", "Top-k: 5 documents", "Prompt template: safety-oriented"],
    "Independent Variable",
    ["Pipeline 1: Vanilla LLM (no retrieval)", "Pipeline 2: Standard RAG (dense semantic)", "Pipeline 3: Multi-Query Expansion", "Pipeline 4: Hybrid Retrieval (RRF fusion)", "Pipeline 5: Query Reformulation + Reranking"]
)

# Slide 12: Evaluation Metrics
add_two_column_slide(prs, "Evaluation Framework",
    "RAGAS Metrics (Grounding Quality)",
    ["Faithfulness — claims supported by context", "Answer Relevance — addresses question intent", "Context Precision — signal-to-noise ratio", "Context Recall — completeness of evidence"],
    "DeepEval Metrics (Safety & Reliability)",
    ["Hallucination Rate — proportion of unsupported claims", "Groundedness — alignment with retrieved evidence", "Correctness — factual accuracy vs reference", "Safety Compliance — adherence to clinical guidelines"]
)

# Slide 13: Section — Results
add_section_slide(prs, 4, "Results & Analysis")

# Slide 14: Faithfulness
add_image_slide(prs, "RAGAS Faithfulness: Hybrid Retrieval Wins", 
    "reports/figures/figure_5_1_faithfulness.png",
    "Hybrid retrieval achieves 0.89 faithfulness — a 47-point improvement over the vanilla baseline (0.42)")

# Slide 15: Hallucination
add_image_slide(prs, "DeepEval Hallucination Rate: 81% Reduction Achieved",
    "reports/figures/figure_5_3_hallucination.png",
    "Hybrid retrieval drops hallucination from 0.58 (baseline) to 0.11 — an 81% relative improvement")

# Slide 16: Context Quality
add_image_slide(prs, "Context Precision & Recall",
    "reports/figures/figure_5_2_precision_recall.png",
    "Hybrid retrieval achieves 0.86 precision and 0.83 recall — the best balance across all pipelines")

# Slide 17: Safety
add_image_slide(prs, "Safety Compliance: All RAG Pipelines Exceed Clinical Threshold",
    "reports/figures/figure_5_6_safety.png",
    "The 0.80 clinical safety threshold is exceeded by all RAG pipelines; hybrid reaches 0.92")

# Slide 18: Per-Question-Type
add_image_slide(prs, "Faithfulness by Question Type",
    "reports/figures/figure_5_5_per_type.png",
    "All pipelines perform best on symptom questions and worst on treatment questions")

# Slide 19: Failure Patterns
add_image_slide(prs, "Failure Pattern Distribution",
    "reports/figures/figure_5_7_failures.png",
    "Hybrid retrieval minimizes all failure categories; vanilla baseline shows even distribution")

# Slide 20: Latency
add_image_slide(prs, "Latency vs Accuracy Trade-off",
    "reports/figures/figure_5_8_latency.png",
    "Hybrid retrieval (1,680ms) offers the best accuracy-to-latency ratio for real-time clinical use")

# Slide 21: Correlation
add_image_slide(prs, "Faithfulness vs Hallucination: Strong Negative Correlation",
    "reports/figures/figure_5_9_correlation.png",
    "Pearson r = -0.97 — faithfulness is a reliable proxy for hallucination risk")

# Slide 22: Section — Key Findings
add_section_slide(prs, 5, "Key Findings & Implications")

# Slide 23: Findings
add_content_slide(prs, "Key Findings", [
    "Retrieval grounding DRAMATICALLY reduces hallucination — 62-81% relative improvement",
    "Hybrid retrieval is the standout performer: highest faithfulness (0.89), lowest hallucination (0.11), highest safety (0.92)",
    "Multi-query expansion improves recall but introduces noise — precision-cost trade-off",
    "Query reformulation improves precision but adds significant latency (3,200ms)",
    "All RAG pipelines exceed the clinical safety threshold of 0.80; vanilla baseline scores only 0.45",
    "Symptom questions are easiest; treatment questions remain challenging due to patient-specific factors"
])

# Slide 24: Clinical Implications
add_content_slide(prs, "Implications for Clinical Practice", [
    "Vanilla LLM is UNSAFE for clinical deployment — hallucination rate of 0.58 is unacceptable",
    "Standard RAG provides a good balance: 62% hallucination reduction with only 1,240ms latency",
    "Hybrid retrieval is recommended for high-stakes applications: best safety, acceptable latency",
    "Retrieval grounding naturally promotes cautious language — models see context limitations",
    "Failure pattern analysis directs future research toward consensus-based retrieval mechanisms"
], "The goal is not to replace clinicians but to augment their expertise with trustworthy, evidence-grounded tools.")

# Slide 25: Contributions
add_content_slide(prs, "Contributions to Knowledge", [
    "First controlled comparative benchmark isolating retrieval strategy effects on hallucination reduction",
    "Empirical evidence that hybrid retrieval outperforms single-method retrieval for medical QA",
    "Per-question-type analysis revealing that medication questions benefit most from hybrid retrieval",
    "Structured failure pattern classification directing future research priorities",
    "Integration of safety compliance evaluation into comparative RAG study — addressing a critical gap"
])

# Slide 26: Section — Conclusion
add_section_slide(prs, 6, "Conclusion & Future Work")

# Slide 27: Limitations
add_content_slide(prs, "Limitations & Mitigations", [
    "Limited to consumer health questions — results may not generalize to clinician-facing content",
    "Automatic metrics correlate with but do not replace clinician judgment",
    "Fixed to GPT-4o-mini — newer models may shift absolute scores but relative rankings should hold",
    "English-only dataset — multilingual medical RAG remains unexplored",
    "No interaction effects tested — optimal top-k may depend on chunk size"
])

# Slide 28: Future Work
add_content_slide(prs, "Future Directions", [
    "Extend hybrid retrieval to include knowledge graphs for drug interaction reasoning",
    "Implement time-aware retrieval to handle newly approved drugs and revised guidelines",
    "Validate automatic metrics against clinician ratings for external validity",
    "Extend evaluation to multilingual medical datasets",
    "Develop granular safety taxonomy covering drug interactions, contraindications, and emergency recognition"
])

# Slide 29: Thank You
add_title_slide(prs,
    "Thank You",
    "Questions & Discussion")

# ==================== SAVE ====================
output_path = r'C:\Users\HP\Downloads\LJMU_Thesis_Defence_Presentation.pptx'
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
